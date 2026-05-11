"""Thesis defense PPTX v5 - based on user's theme template."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os, copy

TEMPLATE = "主题模板.pptx"
FIGS = "figures"

# Load template as base - preserves theme colors, fonts, backgrounds
prs = Presentation(TEMPLATE)
# Set 16:9 widescreen
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# Delete all template slides (keep layouts/theme)
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

# Theme-compatible colors (derived from template feel)
NAVY   = RGBColor(0x1B,0x3A,0x5C)
BLUE   = RGBColor(0x2B,0x57,0x9A)
GREEN  = RGBColor(0x1A,0x7A,0x3A)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
DARK   = RGBColor(0x2D,0x2D,0x2D)
GRAY   = RGBColor(0x77,0x77,0x77)
LIGHT  = RGBColor(0xF5,0xF7,0xFA)
ACCENT = RGBColor(0xEE,0xF2,0xF9)

TITLE_TEXT = "基于多源数据融合的智能天气分析系统设计与实现"
BLANK_LAYOUT = prs.slide_layouts[6]  # blank

def ns():
    return prs.slides.add_slide(BLANK_LAYOUT)

def tbar(slide, text):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.85))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    tf = bar.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.5); tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = WHITE

def tx(slide, l, t, w, h, text, sz=Pt(15), cl=DARK, b=False, al=PP_ALIGN.LEFT):
    bx = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = bx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = sz; p.font.color.rgb = cl; p.font.bold = b; p.alignment = al
    return tf

def tblk(slide, l, t, w, h, lines, sz=Pt(14), cl=DARK, sp=Pt(5)):
    bx = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = bx.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(line, tuple):
            prefix, rest = line
            r1 = p.add_run(); r1.text = prefix; r1.font.size = sz; r1.font.bold = True; r1.font.color.rgb = NAVY
            r2 = p.add_run(); r2.text = rest; r2.font.size = sz; r2.font.color.rgb = cl
        else:
            p.text = line; p.font.size = sz; p.font.color.rgb = cl
        p.space_after = sp

def img(slide, path, l, t, w=None, h=None):
    fp = os.path.join(FIGS, path)
    if not os.path.exists(fp): return None
    if w and h: return slide.shapes.add_picture(fp, Inches(l), Inches(t), Inches(w), Inches(h))
    if w: return slide.shapes.add_picture(fp, Inches(l), Inches(t), width=Inches(w))
    if h: return slide.shapes.add_picture(fp, Inches(l), Inches(t), height=Inches(h))
    return slide.shapes.add_picture(fp, Inches(l), Inches(t))

def card(slide, l, t, w, h, fill=LIGHT, brdr=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = brdr if brdr else fill; shape.line.width = Pt(1)
    # Add subtle shadow via XML
    shape.shadow.inherit = False
    return shape

def callout(slide, l, t, w, h, text, sz=Pt(12), cl=NAVY):
    card(slide, l, t, w, h, fill=ACCENT, brdr=BLUE)
    tx(slide, l+0.15, t+0.06, w-0.3, h-0.12, text, sz=sz, cl=cl)

def pn(slide, n):
    tx(slide, 12.3, 7.12, 0.9, 0.3, str(n), sz=Pt(10), cl=RGBColor(0xAA,0xAA,0xAA), al=PP_ALIGN.RIGHT)

def nt(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def mktable(slide, l, t, w, h, hdrs, rows, hl=None):
    nr = len(rows)+1; nc = len(hdrs)
    ts = slide.shapes.add_table(nr, nc, Inches(l), Inches(t), Inches(w), Inches(h))
    tt = ts.table
    for c, hdr in enumerate(hdrs):
        cell = tt.cell(0, c); cell.text = hdr
        for p in cell.text_frame.paragraphs: p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tt.cell(r+1, c); cell.text = str(val)
            for p in cell.text_frame.paragraphs: p.font.size = Pt(15); p.alignment = PP_ALIGN.CENTER
            is_hl = hl and (r,c) in hl
            for p in cell.text_frame.paragraphs: p.font.bold = is_hl; p.font.color.rgb = GREEN if is_hl else DARK
            cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT if r%2==0 else WHITE

def sec_div(slide, num, title, desc=""):
    tx(slide, 0.8, 1.8, 3.0, 2.0, str(num).zfill(2), sz=Pt(88), cl=RGBColor(0xDD,0xE5,0xF0), b=True)
    tx(slide, 1.2, 3.8, 10.0, 0.8, title, sz=Pt(34), cl=NAVY, b=True)
    if desc: tx(slide, 1.2, 4.7, 10.0, 0.8, desc, sz=Pt(15), cl=GRAY)
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(5.5), Inches(2.5), Pt(3))
    ln.fill.solid(); ln.fill.fore_color.rgb = BLUE; ln.line.fill.background()

def step_card(slide, num, title, desc, l, t, w, h):
    card(slide, l, t, w, h, fill=WHITE, brdr=BLUE)
    tx(slide, l+0.12, t+0.06, 0.35, 0.3, str(num), sz=Pt(20), cl=BLUE, b=True)
    tx(slide, l+0.5, t+0.06, w-0.65, 0.3, title, sz=Pt(13), cl=NAVY, b=True)
    tx(slide, l+0.12, t+0.48, w-0.24, h-0.55, desc, sz=Pt(10.5), cl=DARK)

print(f"Building PPTX v5 based on template '{TEMPLATE}'...")

# ════════════ SLIDE 1: COVER ════════════
s=ns()
top=s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(2.8))
top.fill.solid(); top.fill.fore_color.rgb = NAVY; top.line.fill.background()
img(s,"school_logo.png",5.8,0.15,h=0.9)
tx(s,1.5,1.1,10.3,0.5,"山西财经大学  本科毕业论文答辩",sz=Pt(18),cl=RGBColor(0xBB,0xCC,0xDD),al=PP_ALIGN.CENTER)
tx(s,1.5,1.6,10.3,1.1,TITLE_TEXT,sz=Pt(32),cl=WHITE,b=True,al=PP_ALIGN.CENTER)
ln=s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.2), Inches(2.8), Inches(4.9), Pt(2))
ln.fill.solid(); ln.fill.fore_color.rgb = BLUE; ln.line.fill.background()
tx(s,1.5,3.3,10.3,0.5,"答辩人: 徐毅成    学号: 202207020136    专业: 信息管理与信息系统",sz=Pt(16),cl=DARK,al=PP_ALIGN.CENTER)
tx(s,1.5,4.1,10.3,0.5,"指导教师",sz=Pt(15),cl=GRAY,al=PP_ALIGN.CENTER)
for i,(num,label) in enumerate([("5","核心功能模块"),("25+","API接口"),("~8000","行Python代码")]):
    tx(s,2.5+i*3.0,5.2,2.5,0.5,num,sz=Pt(28),cl=NAVY,b=True,al=PP_ALIGN.CENTER)
    tx(s,2.5+i*3.0,5.7,2.5,0.4,label,sz=Pt(12),cl=GRAY,al=PP_ALIGN.CENTER)
nt(s,"各位老师好，我是徐毅成。我的论文题目是基于多源数据融合的智能天气分析系统设计与实现。")

# ════════════ SLIDE 2: OUTLINE ════════════
s=ns();tbar(s,"汇报提纲")
tblk(s,0.8,1.5,7.5,5.0,[
    ("一、","研究背景与意义 —— 为什么要做这个系统"),
    ("二、","系统总体设计 —— 整体架构和数据流转"),
    ("三、","五大核心功能模块 —— 重点部分，展示实际页面截图"),
    ("四、","系统测试与结果分析 —— 关键逻辑和ML效果数据"),
    ("五、","总结与展望 —— 做了什么、有什么不足、以后怎么改进"),
],sz=Pt(18))
tx(s,0.8,6.5,10.0,0.4,"答辩时长约15分钟  |  开发周期约4个月  |  全部独立设计、开发、撰写",sz=Pt(12),cl=GRAY)
pn(s,2); nt(s,"先让老师对整个汇报结构有预期。")

# ════════════ SLIDE 3: BACKGROUND ════════════
s=ns();tbar(s,"研究背景: 我们为什么需要一个更智能的天气平台?")
tblk(s,0.5,1.1,7.8,2.8,[
    ("先看现状: ","大多数人查天气就看一个App——告诉你明天多少度、会不会下雨。但如果你想知道这个预报靠不靠谱? 同一时间不同渠道给出的温度能差好几度，到底信哪个? 农业生产者不光要知道下不下雨，还想知道这场雨对正在抽穗的水稻有多大影响；气象爱好者想了解今天大气稳不稳定。这些需求，普通天气App满足不了。"),
    ("数据其实都有: ","欧洲中期天气预报中心(ECMWF)的全球预报、美国(GFS)和德国(ICON)的模式结果、怀俄明大学的探空气球数据、ERA5-Land几十年的再分析气候资料——全是公开免费的。但散落在不同网站，格式各异，时间口径不统一，普通人没法直接用。杭州缺一个把这些整合起来的综合平台。"),
],sz=Pt(13.5))
tx(s,0.5,4.1,12.3,0.4,"这个系统的本质: 用代码重现一个专业气象分析人员的完整思考链条",sz=Pt(13),cl=BLUE,b=True)
# Horizontal 5-step flow
steps=[("1 看未来","多模式预报\n回答会发生什么"),
       ("2 看当前","探空分析\n回答为什么发生"),
       ("3 看过去","历史气候\n提供历史参照系"),
       ("4 修正偏差","机器学习\n让预报结果更可靠"),
       ("5 看影响","农业气象\n转化成该怎么办")]
for i,(title,desc) in enumerate(steps):
    x=0.5+i*2.5
    # Arrow between boxes (except after last)
    if i<4:
        arr=s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,Inches(x+2.1),Inches(5.05),Inches(0.35),Inches(0.18))
        arr.fill.solid();arr.fill.fore_color.rgb=BLUE;arr.line.fill.background()
    # Card
    card(s,x,4.6,2.1,0.9,fill=WHITE,brdr=BLUE)
    tx(s,x+0.08,4.65,1.95,0.3,title,sz=Pt(11),cl=NAVY,b=True)
    tx(s,x+0.08,4.95,1.95,0.5,desc,sz=Pt(9.5),cl=DARK)
pn(s,3); nt(s,"用聊天语气讲背景，引入横排五步思考链条——不是堆功能，是模拟专业分析流程。")

# ════════════ SLIDE 4: CONTRIBUTIONS ════════════
s=ns();tbar(s,"我做了哪些工作")
tblk(s,0.5,1.2,11.5,6.0,[
    ("1. 多模式预报对比: ","同时接入ECMWF(欧洲)、GFS(美国)、ICON(德国)三种全球数值预报。72小时逐小时曲线同轴叠加，附带75%置信区间集合预报，把不确定性量化出来。"),
    ("2. 探空专业诊断: ","从怀俄明大学抓取探空数据，画Skew-T图，算CAPE/CIN/K指数等参数，翻译成通俗风险等级。自己设计了一张卡通大气体检图。"),
    ("3. EW指数极端事件识别: ","自己设计的方法——把每一天的温/雨值放到过去30年同月同日样本里比，看它到底有多罕见。经过过程去重和四季分组，给出有统计意义的年度代表事件榜单。"),
    ("4. ML误差校正: ","ECMWF落到杭州站点有系统性偏差。用97个历史预报文件+对应真实观测训练随机森林模型。温度误差降22%，风速降40%。证明口径统一比换复杂模型重要。"),
    ("5. 农业气象服务: ","浙江5种典型作物，每种建立了生长阶段/温度阈值/灾害风险知识库。系统逐天扫描未来7天预报并给出农事建议。"),
    ("6. 工程完整性: ","五层架构、统一站点口径、训练可追溯、多处容错回退。"),
],sz=Pt(12.5))
pn(s,4); nt(s,"六项工作快速过。每项有具体数据。")

print("Slides 1-4 done")

# ════════════ SLIDE 5: SECTION DIVIDER ════════════
s=ns();sec_div(s,1,"系统总体设计","五层架构 · 数据流转 · 站点口径统一"); pn(s,5)

# ════════════ SLIDE 6: ARCHITECTURE ════════════
s=ns();tbar(s,"系统是怎么搭起来的: 五层架构")
img(s,"system_architecture.png",0.2,0.9,w=12.9)
for i,(lb,sb) in enumerate(zip(
    ["展示层","路由与接口层","服务层","分析层","模型层"],
    ["12个HTML页面+Bootstrap\n用户看到的网页界面","app.py管理25+路由\n页面跳转和API调度","数据抓取/缓存/转换\n外部API数据统一格式","气象参数/统计/农业\n分析逻辑都在这一层","ML校正模型加载调用\n随机森林+梯度提升"]
)):
    tx(s,0.6+i*2.55,5.1,2.35,0.42,lb,sz=Pt(12),cl=NAVY,b=True,al=PP_ALIGN.CENTER)
    tx(s,0.6+i*2.55,5.5,2.35,0.8,sb,sz=Pt(8.5),cl=GRAY,al=PP_ALIGN.CENTER)
callout(s,0.5,6.4,12.3,0.45,"核心设计: 所有模块统一锚定杭州58457站(30.24N/120.15E)，时间统一按UTC标准时次对齐。ML校正可信度的根基。",sz=Pt(12))
pn(s,6); nt(s,"五层架构，强调统一站点口径。")

# ════════════ SLIDE 7: DATA FLOW ════════════
s=ns();tbar(s,"数据是怎么跑起来的: 从外部到屏幕的完整路径")
img(s,"data_flow.png",0.2,1.1,h=5.0)
tblk(s,6.8,1.2,6.0,5.8,[
    ("四类外部数据源(全部公开免费): ","(1)Open-Meteo——统一提供ECMWF/GFS/ICON/GEM四种模式逐小时预报；(2)怀俄明大学探空资料库——全球探空站每12小时放气球记录温湿风廓线；(3)RP5历史站点观测——杭州地面站多年逐日记录；(4)ERA5-Land再分析——欧洲中心回算的几十年格点气候数据，时间连续性好。"),
    ("数据清洗和口径统一: ","空间上统一到杭州58457站(同一经纬度和海拔)；时间上把北京时转成UTC标准时次(00Z/12Z)。如果ML训练用A站、运行用B站，校正结果没法解释。"),
    ("各模块按需分析: ","清洗好的数据分发到探空(参数/风险评估)、历史(统计/趋势/EW)、农业(积温/适宜度/风险扫描)，处理完通过API传给前端展示。"),
],sz=Pt(12))
pn(s,7); nt(s,"四类数据源各用一句话说清楚。")

# ════════════ SLIDE 8: HOMEPAGE ════════════
s=ns();tbar(s,"系统首页: 一个页面看清全局")
img(s,"homepage_full.png",0.2,1.0,w=12.9)
img(s,"home_navigation_flow.png",9.8,6.95,w=3.2)
tblk(s,0.2,6.9,9.5,0.45,[("顶部: ","实时天气+7天概览 | 中部: 功能模块导航卡片 | 右侧: 智能提醒 | 右下: 首页导航流程设计")],sz=Pt(11))
pn(s,8); nt(s,"首页三个角色: 实时感知、模块分发、智能提醒。")

print("Slides 5-8 done")

# ════════════ SLIDE 9: FORECAST - CONCEPT ════════════
s=ns();tbar(s,"多模式预报(一): 不只告诉你多少度,还能告诉你可以信几分")
img(s,"advanced-forecast1.png",7.9,1.1,w=5.1)
tx(s,7.9,6.5,5.1,0.3,"72小时逐小时温度/降水/湿度/风速曲线(Plotly交互图)",sz=Pt(9.5),cl=GRAY,al=PP_ALIGN.CENTER)
tblk(s,0.5,1.2,7.2,5.8,[
    ("一个朴素的问题: ","同一个时间、同一个地点，ECMWF预报25度，GFS预报23度，ICON预报26度——到底信哪个? 普通App只给你一个数，不告诉你背后可能有2-3度的分歧。"),
    ("我的做法: ","从Open-Meteo同时拿三种全球模式的72小时逐小时预报，不同颜色曲线叠在同一张图里。三条线紧贴=可信度高；分得开=要留个心眼。还做了集合预报——四模式一起算均值和中位数，画25%-75%灰色置信区间带。"),
    ("四个模式: ","ECMWF(欧洲)——国际公认综合性能最好，ML校正主链路；GFS(美国)——免费开放更新频繁；ICON(德国)——较新的非静力模式；GEM(加拿大)——偶尔有惊喜。"),
    ("技术细节: ","四模式Open-Meteo一个API并行获取，Plotly.js交互曲线(悬停看数值)，页面异步加载(首屏不白屏)，5分钟短时缓存。"),
],sz=Pt(12.5))
pn(s,9); nt(s,"讲清四个模式和集合预报的置信区间。")

# ════════════ SLIDE 10: FORECAST - DEMO ════════════
s=ns();tbar(s,"多模式预报(二): 实际页面效果")
img(s,"advanced-forecast2.png",0.2,1.1,w=6.5)
tx(s,0.2,6.5,6.5,0.3,"多模式同轴对比: ECMWF(红)/GFS(蓝)/ICON(橙)+灰色置信区间",sz=Pt(10),cl=GRAY,al=PP_ALIGN.CENTER)
img(s,"advanced-forecast3.png",7.0,1.1,w=6.1)
tx(s,7.0,6.2,6.1,0.3,"7天趋势+AI智能分析面板(自动生成文字版天气解读)",sz=Pt(10),cl=GRAY,al=PP_ALIGN.CENTER)
tblk(s,0.3,6.9,12.7,0.55,[
    ("左: ","三模式同轴对比页——哪个偏暖、哪个偏冷、谁跟谁一致，一望而知。"),
    ("右: ","7天趋势+AI——每日最高最低温+降水概率+置信度标签(绿=高/黄=中/红=低)+AI自然语言天气解读。"),
],sz=Pt(11.5))
pn(s,10); nt(s,"指着截图讲解。")

print("Slides 9-10 done")

# ════════════ SLIDE 11: UPPER AIR - CONCEPT ════════════
s=ns();tbar(s,"探空分析(一): 天气预报看地面,探空看高空——好比大气的CT扫描")
img(s,"upperair_flow.png",0.2,1.1,h=5.3)
tblk(s,4.8,1.2,8.0,5.8,[
    ("它解决什么问题? ","地面预报告诉你明天可能下大雨。但这场雨是普通降雨还是强对流暴雨? 大气现在是稳定还是不稳定? 这些地面数据看不出来。探空就是放气球到高空，测量每一层的温度、湿度、风速风向——好比给大气做CT扫描。"),
    ("数据从哪来? ","全球约800个探空站每天UTC 00和12点(北京时8点和20点)各放一次气球。从美国怀俄明大学公开资料库自动抓取杭州站(58457)数据。"),
    ("拿回来怎么处理? ","解析(原始文本转结构化数据)->计算(MetPy算CAPE/CIN/K指数/抬升指数/风切变等十几个参数)->画图(Skew-T斜温图)->评估(根据CAPE和风切变判断强对流风险等级)。"),
    ("时间换算: ","页面输入北京时->后端自动转UTC->找最近探空时次。如果数据还没出来自动回退12小时，返回结果标注实际用了哪个时次。"),
],sz=Pt(11.5))
pn(s,11); nt(s,"CT扫描比喻。完整流程讲解。")

# ════════════════ SLIDE 12: UPPER AIR - GRAPHICS + CARTOON ════════════════
s=ns();tbar(s,"探空分析(二): Skew-T图+卡通大气体检图——专业和通俗两不误")
img(s,"upperair2.png",0.2,1.1,w=6.2)
tx(s,0.2,6.4,6.2,0.3,"Skew-T图: 红=温度廓线, 蓝=露点廓线, 右侧=风矢量, 紫色区=CAPE",sz=Pt(9.5),cl=GRAY,al=PP_ALIGN.CENTER)
img(s,"upperair_cartoon.png",6.7,1.1,w=6.3)
tx(s,6.7,6.4,6.3,0.3,"卡通大气体检图(自己设计): 温度计/风速仪/湿度表直观展示各层大气状态",sz=Pt(9.5),cl=GRAY,al=PP_ALIGN.CENTER)
tblk(s,0.2,6.75,12.8,0.7,[
    ("左: ","Skew-T标准探空图——红蓝线距离代表饱和程度，越近越易成云致雨。紫色CAPE是大气中储存的燃料。"),
    ("右: ","卡通直观图——自己设计实现。把专业探空数据翻译成温度计、风速仪、湿度表这些谁都看得懂的视觉元素。四层(高空/中层/0度层/地面)，绿=正常/橙=注意/红=异常。即使完全不懂气象也能一眼看出今天大气健不健康。"),
],sz=Pt(11.5))
pn(s,12); nt(s,"两张图对比——左给专业，右给普通。卡通图是自己设计实现的。")

# ════════════════ SLIDE 13: UPPER AIR - PARAMETERS ════════════════
s=ns();tbar(s,"探空分析(三): 参数解读与强对流风险评估")
img(s,"upperair1.png",0.2,1.1,w=7.3)
tx(s,0.2,6.6,7.3,0.3,"探空参数面板: 关键指数+中文说明+风险等级卡片",sz=Pt(10),cl=GRAY,al=PP_ALIGN.CENTER)
tblk(s,7.8,1.2,5.2,6.0,[
    ("核心参数(都配有通俗说明):",""),
    ("CAPE(对流有效位能): ","大气不稳定能量的量化。>1000=可能打雷；>2500=雷暴大风冰雹。"),
    ("CIN(对流抑制): ","阻止空气上升的盖子。CIN大=即使CAPE高也不易触发。"),
    ("K指数: ","综合温湿度指标。>35=雷暴可能性增大。"),
    ("0-6km风切变: ",">12m/s配合高CAPE=可能超级单体风暴。"),
    ("风险等级: ","低(绿)=CAPE<500且K<35；中(橙)=CAPE>1000或K>35且切变>12；高(红)=CAPE>1000且切变>12。"),
    ("三视角解读: ","航空(风切变颠簸积冰)/农业(强对流冰雹)/公众(户外安全)。"),
    ("异常处理: ","CAPE为负数->统一标N/A，不显示可能误导人的数字。"),
],sz=Pt(11))
pn(s,13); nt(s,"参数业务含义和风险判定规则。")

print("Slides 11-13 done")

# ════════════════ SLIDE 14: HISTORY - OVERVIEW ════════════════
s=ns();tbar(s,"历史气候分析(一): 从几十年数据里看出规律")
img(s,"history.png",0.2,1.1,w=7.5)
tx(s,0.2,6.6,7.5,0.3,"年度气候概览页: 统计卡片+月度对比图+极端事件榜单",sz=Pt(10),cl=GRAY,al=PP_ALIGN.CENTER)
tblk(s,8.0,1.2,5.0,6.0,[
    ("历史气候模块做三件事:",""),
    ("第一,年度概览: ","给定某一年，自动统计年均温/年降水/雨雪日/高温日(>=35)/低温日(<=0)/暴雨日(>=50mm)，跟标准气候态(1981-2010三十年均值)对比——偏暖还是偏冷、偏湿还是偏干。"),
    ("第二,长期趋势: ","把多年统计结果串起来，看气温/降水/极端频次的变化趋势，算出十年变化速率。"),
    ("第三,找年度代表极端事件: ","传统做法拉排行榜——最高温哪天、最大雨哪天。但绝对值高不等于历史异常。EW指数做了改进: 把每一天的值放到过去30年同月同日样本里比，看排第几。排越前说明真的不寻常。"),
    ("数据从哪来? ","杭州站RP5格式多年逐日观测资料，近30个字段。RP5降水是滚动累积的，需要复杂层级解析还原成日降水量。"),
],sz=Pt(11))
pn(s,14); nt(s,"历史气候三件事。引出EW指数动机。")

# ════════════════ SLIDE 15: EW INDEX ════════════════
s=ns();tbar(s,"历史气候分析(二): EW指数——怎么才算真正异常的天气?")
img(s,"history_ew_flow.png",0.2,1.1,h=5.2)
tblk(s,4.3,1.2,8.5,6.0,[
    ("举个例: ","7月15日杭州最高温38度。但过去30年每年这天平均37.5度——那只高0.5度，排60%位置，挺正常的。3月2日最高温22度虽不高，但历史同日平均才15度——排95%，这才是真正异常的暖事件。"),
    ("EW指数就是做了这件事: ","把每一天的高温/低温/降水值放回它自己这个日期(月+日)的历史样本池，算它排第几。排位越高说明这一年这天在历史中越罕见。"),
    ("四个设计细节: ","(a)低温反向——越冷分越高；(b)同季门槛——百分位高但绝对值不够极端的排除；(c)过程去重——连5天高温只取最极端1天；(d)四季分组——春夏秋冬各出各榜，每季每要素最多3条，避免夏天占满全部名额。"),
    ("趋势页: ","系统还有一个趋势页——把多年统计串起来看气温/降水/极端频次的年际变化，每条趋势线标注十年变化速率。比如2015到2025年，杭州年均温以约0.3度/十年的速度上升。"),
],sz=Pt(11.5))
pn(s,15); nt(s,"用具体例子讲透EW逻辑。趋势页内容也合并到这里。")

# ════════════════ SLIDE 16: HISTORY - EW PAGE ════════════════
s=ns();tbar(s,"历史气候分析(三): EW指数页面——四季分组代表事件榜单")
img(s,"ew_page.png",0.3,1.1,h=5.8)
tx(s,0.3,6.9,3.0,0.3,"EW指数代表事件榜单: 四季分组，可切换原始极值/EW模式",sz=Pt(9.5),cl=GRAY,al=PP_ALIGN.CENTER)
tblk(s,3.6,1.2,9.3,5.5,[
    ("页面交互: ","默认显示传统原始极值榜。右上角一键切到EW指数模式——榜单变成四季分组(春/夏/秋/冬)，每季每种要素(高温/低温/降水)最多各3条代表事件。"),
    ("为什么四季分组? ","不同季节天气背景完全不同。夏天极端高温和冬天极端低温没法放一起比。分组后各自评选才合理。每季限3条避免某个季节(如夏天)占满全部名额。"),
    ("双气候态: ","支持1981-2010和1991-2020两套标准气候态切换。气候均值本身在缓慢上升，同一年的距平用不同基准算出来不一样。WMO气候业务标准做法。"),
    ("两种榜单对比: ","传统极值榜选绝对值最大的，EW榜选历史背景中最异常的。前者告诉你天气的绝对边界，后者告诉你哪些过程真正值得记住。"),
],sz=Pt(13))
pn(s,16); nt(s,"左侧竖长截图缩小适配，右侧详细解释。")

print("Slides 14-16 done")

# ════════════════ SLIDE 17: ML - CONCEPT ════════════════
s=ns();tbar(s,"ML误差校正(一): 用机器学习让ECMWF的预报更贴近杭州实际")
img(s,"ml_forecast_apply.png",7.3,1.2,w=5.8)
tx(s,7.3,5.5,5.8,0.3,"多模式预报页中ML校正效果: 红色(原始ECMWF) vs 蓝色(ML校正后)",sz=Pt(10),cl=GRAY,al=PP_ALIGN.CENTER)
tblk(s,0.5,1.2,6.6,6.2,[
    ("打个比方: ","ECMWF好比世界级大厨——通用菜谱到了杭州厨房总是偏咸(偏暖1-2度)、偏油(偏湿)。我要做的不是质疑大厨，而是训练一个本地助手，告诉他当大厨说25度时杭州人尝到的大概是23度。"),
    ("具体怎么做: ","收集新版ECMWF预报文件97个(每个含18个逐小时参数)，用杭州58457站同时段真实观测(Meteostat)做标准答案。按起报时间对齐后训练随机森林模型。温度/湿度/风速各一个回归模型，降水拆两步——先分类(下不下)、再回归(下多少)。"),
    ("最关键的不是模型: ","花最多精力的不是调参数，而是把站点口径和时次口径彻底统一。旧版训练用A站、运行用B站——结果没法解释。新版全链路统一到杭州58457。时次推断: 凌晨4:15采集->前日12Z，下午4:15->当日00Z。训练按起报时间切分——严禁随机打散导致信息泄漏。"),
    ("右图展示效果: ","多模式预报页面中，ECMWF原始曲线(红)和ML校正后(蓝)叠在一起。一键切换，模型不可用时自动回退原始值。"),
],sz=Pt(12))
pn(s,17); nt(s,"大师傅+本地助手比喻。强调口径统一。")

# ════════════════ SLIDE 18: ML - TRAINING ════════════════
s=ns();tbar(s,"ML误差校正(二): 训练流程——从97个文件到5个模型")
img(s,"ml_correction_flow.png",0.2,1.1,h=5.5)
tblk(s,4.3,1.2,8.5,6.0,[
    ("训练数据审计: ","100个文件->97个通过->3个因采集时间不符ECMWF起报窗口剔除。训练集2026.2.15~4.3，测试集2026.4.4~4.10。虽只2个月数据，但统一口径下已能看到明确校正效果。"),
    ("V1模型(当前方案): ","scikit-learn随机森林。温度/湿度/风速各一个回归模型；降水两步——先分类(>0.1mm为有雨)，再回归(估量)。总共5个模型文件。"),
    ("特征设计(12维): ","原始预报值+时间特征(小时/月份/工作日/周末)+预报时效lead_hours+时效分桶+循环编码(小时正弦余弦->23点和0点在数值上接近)。V2扩展到24维: 加季节/起报周期/3h-6h变化率。"),
    ("为什么降水两步法? ","降水分两个问题: 下不下(分类)和下多少(回归)。混在一起模型可能为了拟合量级牺牲对有无的判断。拆开后分类指标虽没全面提升，但量级估计明显更准——对农业来说知道具体下多少比知道会不会下更重要。"),
],sz=Pt(11.5))
pn(s,18); nt(s,"训练流程完整讲解: 审计->特征->模型选择->降水两步法。")

# ════════════════ SLIDE 19: ML - RESULTS ════════════════
s=ns();tbar(s,"ML误差校正(三): 数据说话——三个变量都有改善,风速效果最明显")
mktable(s,2.0,1.3,9.0,2.2,
    ["气象要素","原始ECMWF MAE","校正后 MAE","改进幅度"],
    [["温度","2.355度","1.832度","降低22.2%"],
     ["湿度","9.851%","9.085%","降低7.8%"],
     ["风速","4.413 m/s","2.613 m/s","降低40.8%"]],
    hl={(0,3),(2,3)})
tblk(s,1.0,3.7,11.3,2.8,[
    ("解读: ","(1)风速改善最明显(降40%)——原始ECMWF在杭州站风速偏差大，校正空间大；(2)短时效(1-3天)改善好于长时效；(3)湿度改善相对小——局地影响大、变化快。"),
    ("降水诚实说: ","分类指标与原始持平，没有全面改善。但雨样本MAE和全样本MAE明显降低——判定有雨后量级估计比以前准。策略偏保守(减少误报)，召回率略降、精确率不变或略升。量级改善在农业灌溉、水库调度、排涝准备中更有实际意义。"),
    ("评估条件: ","严格时间隔离测试集(4.4-4.10)，3折交叉验证，指标记录在metrics_report.json中可回溯复现。"),
],sz=Pt(12.5))
pn(s,19); nt(s,"核心数据页。")

# ════════════════ SLIDE 20: ML - V1 VS V2 ════════════════
s=ns();tbar(s,"ML误差校正(四): V1(随机森林) vs V2(CatBoost)——诚实的结果")
mktable(s,1.5,1.3,10.0,2.4,
    ["指标","原始ECMWF","V1(RandomForest)","V2(CatBoost)"],
    [["温度 MAE","2.355","1.832","2.220"],
     ["湿度 MAE","9.851","9.085","10.076"],
     ["风速 MAE","4.413","2.613","3.061"],
     ["降水 Accuracy","0.754","0.743","0.658"]],
    hl={(0,2),(1,2),(2,2)})
tblk(s,1.0,4.0,11.3,3.0,[
    ("V2(CatBoost)未超过V1(RandomForest)——合理且诚实的结果。",""),
    ("为什么? ","当前约2个月样本(97个文件)。小样本下随机森林(bagging)通过Bootstrap天然比CatBoost(boosting)更稳定。CatBoost需更多样本和更细超参搜索才能发挥优势。"),
    ("这个发现本身有价值: ","说明在ECMWF校正任务中，当前最关键的不是换复杂模型，而是控数据质量——站点统一、时次一致、按起报时间切分防泄漏。后续积攒整年数据后重新评估boosting才公平。"),
],sz=Pt(12))
pn(s,20); nt(s,"V2没赢V1是诚实发现——口径统一比换模型重要。")

print("Slides 17-20 done")

# ════════════════ SLIDE 21: AGRO - DESIGN ════════════════
s=ns();tbar(s,"农业气象(一): 把天气预报转化成田间管理建议")
img(s,"agro_service_flow.png",0.2,1.1,h=5.2)
tblk(s,4.3,1.2,8.5,5.8,[
    ("这个模块做什么? ","前面讲的都在分析和展示气象信息。农业模块多走一步——把气象能力嫁接到具体作物管理上。每种作物有温度底线、水分需求、生长节奏和怕的灾害。系统把未来7天预报和作物知识库比对，给出有针对性的农事建议。"),
    ("五5种浙江代表性作物:",""),
    ("水稻: ","4月中播种，全生育期需积温2200度(基准10度)。最怕抽穗扬花期(9月上旬)寒露风(<17度灌浆受阻)和高温热害(>35度花粉失活)。"),
    ("西湖龙井: ","2月中萌芽。最怕倒春寒——萌芽期<4度持续3天嫩芽冻伤。7-8月防高温灼伤(>35度连续5天叶片焦枯)。"),
    ("杨梅: ","6月初-7月初成熟采摘。最怕暴雨(>20mm)落果、大风(>10m/s)吹落、高湿(>90%)烂果。浙江梅雨季撞上成熟期，风险极高。"),
    ("柑橘: ","全年生长期长。最怕冬季极端低温(<-2度冻害死树)和夏季日灼(>38度灼伤果皮)。"),
    ("小白菜: ","全年可种周期短(~50天)。怕高温(>32度徒长)和暴雨(>50mm渍涝)。积温起算日由用户自设播种日期。"),
],sz=Pt(11.5))
pn(s,21); nt(s,"五种作物各讲最核心风险点。")

# ════════════════ SLIDE 22: AGRO - DASHBOARD ════════════════
s=ns();tbar(s,"农业气象(二): 仪表盘——总览全局农情")
img(s,"agro-dashboard.png",0.5,1.2,w=12.3)
tblk(s,0.5,6.2,12.3,1.2,[
    ("总览页: ","顶部滚动条逐条列出各作物当前风险(红=高风险/橙=中风险/蓝=阶段提醒/绿=适宜作业)；中部作物卡片显示当前生长阶段、适宜度评分(0-100)、风险标签；底部7天农事日历标注每日风险等级和操作建议；右侧AI农情研判(DeepSeek生成，API不可用自动切规则引擎)。"),
    ("双层结构: ","总览页给管理者快速全局判断；点进任意作物专题页(/crop/<id>)看积温进度、环境曲线、逐日风险清单和个性化建议。合起来就是天气->作物->决策的完整闭环。"),
],sz=Pt(14))
pn(s,22); nt(s,"一张大图占满，下面简洁说明。")

# ════════════════ SLIDE 23: AGRO - CROP MONITOR (图4) ════════════════
s=ns();tbar(s,"农业气象(三): 作物生长环境监测")
img(s,"crop_monitor.png",0.5,1.2,w=12.3)
tblk(s,0.5,6.5,12.3,0.8,[
    ("作物专题页的环境监测: ","点进水箱或龙井的专题页后，能看到这张生长环境监测面板——左侧展示当前生育阶段、积温累计进度(当前值/目标值/百分比)、预估成熟日期；右侧是未来7天环境曲线(每日最高最低温+降水+适宜度变化)，底部逐日列出风险预警(日期/风险类型/触发阈值/当前预报值/建议应对)。"),
    ("业务逻辑: ","系统拿到未来7天逐小时ECMWF预报->按天聚合(最高温/最低温/总降水/均湿/均风)->跟当前作物风险阈值逐条比对->命中阈值则生成预警->同时检查适合作业(喷药/施肥/修剪)的天气窗口。整套逻辑在agro_alert_engine.py中实现，每种作物可配置自己的风险规则。"),
],sz=Pt(14))
pn(s,23); nt(s,"图4展示生长环境监测面板，交代背后业务逻辑。")

# ════════════════ SLIDE 24: AGRO - RICE ════════════════
s=ns();tbar(s,"农业气象(四): 水稻专题——积温进度与风险预警")
img(s,"rice.png",0.5,1.2,w=12.3)
tblk(s,0.5,6.2,12.3,1.2,[
    ("水稻专题页: ","左栏——积温进度条(当前积温/目标2200度/百分比)+预估成熟日期(基于积温达标时间)；中栏——未来7天每日最高最低温和降水曲线(Chart.js)+适宜度评分变化；右栏——逐日风险预警清单。"),
    ("积温(GDD)怎么算? ","逐日累计: (日最高温+日最低温)/2 - 基准温度(水稻10度)，只取正数。积温用热量而非日历来衡量作物发育——达到2200度日意味着水稻走完了从播种到成熟的全过程。历史积温用ERA5-Land再分析(长期连续)，未来7天用ECMWF预报(时效性强)，两源互补。"),
],sz=Pt(14))
pn(s,24); nt(s,"水稻专题页大图。积温公式和双源策略。")

# ════════════════ SLIDE 25: AGRO - TEA + BAYBERRY ════════════════
s=ns();tbar(s,"农业气象(五): 龙井倒春寒预警 + 杨梅暴雨落果预警")
img(s,"tea.png",0.2,1.2,w=6.3)
tx(s,0.2,6.2,6.3,0.3,"龙井茶专题: 积温进度+倒春寒(<4度)预警+高温灼伤(>35度)预警",sz=Pt(10.5),cl=GRAY,al=PP_ALIGN.CENTER)
img(s,"bayberry.png",6.8,1.2,w=6.3)
tx(s,6.8,6.2,6.3,0.3,"杨梅专题: 成熟期暴雨(>20mm)/大风(>10m/s)/高湿(>90%)落果烂果预警",sz=Pt(10.5),cl=GRAY,al=PP_ALIGN.CENTER)
tblk(s,0.2,6.5,12.8,0.9,[
    ("左——西湖龙井: ","3-4月倒春寒(日最低<4度持续3天)->嫩芽冻伤->提示覆盖防霜/烟熏防冻。7-8月连续高温(>35度持续5天)->叶片灼伤->提示灌溉降温遮阴。倒春寒直接关系春茶产量和品质。"),
    ("右——杨梅: ","6月中下旬成熟采摘期遇暴雨+大风+高湿->落果烂果三连预警。浙江梅雨季(6月中-7月初)正好撞上杨梅成熟期。系统建议暴雨前48小时抢收、雨后排水降湿。"),
    ("系统不只报风险: ","同时扫描未来7天有没有适合作业的窗口——暴雨后48小时无雨+风速<5m/s+温度15-30度->适合喷药/施肥。趋利避害并重。"),
],sz=Pt(13))
pn(s,25); nt(s,"两个作物案例+趋利避害并重的设计理念。")

print("Slides 21-25 done")

# ════════════════ SLIDE 26: TESTING ════════════════
s=ns();tbar(s,"系统测试: 验证的不只是页面能打开,更是关键逻辑对不对")
tblk(s,0.5,1.2,8.5,6.0,[
    ("测试四维度:",""),
    ("1. 功能流程: ","五大模块页面均正常加载，API正常响应，业务链路闭环——预报从获取到ML校正全程跑通；探空从时间输入到风险评估全程跑通；农业从预报到预警生成全程跑通。"),
    ("2. 关键逻辑验证: ","EW事件筛选(反向百分位+过程去重+四季分组+数量控制)通过pytest自动化验证。探空北京时->UTC换算正确(4.10 08:00->4.10 00:00)。预报时次推断稳定(04:15->12Z, 16:15->00Z)。"),
    ("3. 异常处理: ","CAPE负值->N/A。外部API超时->缓存回退。ML模型未加载->返回原始ECMWF。DeepSeek挂了->农业AI切规则引擎。"),
    ("4. ML可追溯: ","严格时间隔离，97/3审计，metrics_report.json记录所有指标可回溯复现。"),
    ("其实最难的不是页面本身: ","是时次推断、EW去重、CAPE N/A、口径统一这些容易被忽略的细节。它们决定系统输出的结果到底可不可信。"),
],sz=Pt(12.5))
img(s,"smart_tips.png",9.2,1.5,w=3.8)
tx(s,9.2,5.0,3.8,0.3,"首页智能提醒区(聚合多模块预警)",sz=Pt(10),cl=GRAY,al=PP_ALIGN.CENTER)
pn(s,26); nt(s,"四个维度展开。细节决定可信度。")

# ════════════════ SLIDE 27: SUMMARY ════════════════
s=ns();tbar(s,"总结: 我做了一个什么样的系统?")
tblk(s,0.5,1.3,12.3,5.5,[
    ("一句话: ","围绕杭州地区、以预报分析为主线、整合探空诊断+历史统计+ML误差校正+农业服务的综合天气分析平台。不是功能堆砌，是用统一口径和五层架构组织起来的工程闭环——本质上是专业气象分析人员完整思考链条的代码复现。"),
    ("核心亮点:",""),
    ("多模式对比: ","ECMWF/GFS/ICON三模式同轴+75%置信区间集合预报，把不确定性量化给用户。"),
    ("EW指数: ","自己设计的极端事件识别，从绝对值排名升级到历史异常程度评估。"),
    ("ML校正: ","温度降22%，风速降40%。更重要的——口径统一比换复杂模型重要。"),
    ("农业服务: ","5种作物+积温+风险预警+AI建议，从天气延伸到行业决策。"),
    ("工程完整性: ","五层解耦+容错回退+训练可追溯。97/100文件审计通过。"),
    ("核心价值: ","不在于新算法，在于围绕真实业务把多源数据、分析和ML整合成可运行、可解释、可追溯的系统。"),
],sz=Pt(14))
pn(s,27); nt(s,"提炼每项工作的独特价值。工程闭环和实践价值。")

# ════════════════ SLIDE 28: LIMITATIONS & THANKS ════════════════
s=ns();tbar(s,"当前不足、未来展望与致谢")
card(s,0.5,1.3,5.8,3.8,fill=LIGHT)
tx(s,0.8,1.5,5.3,0.4,"当前不足",sz=Pt(20),cl=NAVY,b=True)
tblk(s,0.8,2.1,5.3,2.8,[
    ("1. 外部依赖: ","数据源可用性不由本地决定，虽有缓存回退但无系统化本地归档。"),
    ("2. 样本有限: ","只采集了2个月(97个)。可证明校正有效，但盖不完完整季节和所有天气类型。"),
    ("3. 单站点: ","以杭州为核心。聚焦有助于做深，迁移需逐个适配站点和作物。"),
    ("4. 部分规则化: ","EW阈值/农业权重/风险等级基于经验和文献，数据积累后可向数据驱动改进。"),
],sz=Pt(12))
card(s,7.0,1.3,5.8,3.8,fill=LIGHT)
tx(s,7.3,1.5,5.3,0.4,"未来方向",sz=Pt(20),cl=NAVY,b=True)
tblk(s,7.3,2.1,5.3,2.8,[
    ("1. 完善数据层: ","建立本地归档和元数据管理，降低外部接口硬依赖。"),
    ("2. 扩大训练样本: ","积攒整年后按季节/时效分组重训，重新评估boosting模型。"),
    ("3. 多站点扩展: ","站点配置参数化，支持快速部署新城市，增加更多作物和病虫害模型。"),
    ("4. 增强AI分析: ","AI+规则双路径扩展到预报摘要、探空解读、年度报告等更多场景。"),
],sz=Pt(12))
ln2=s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(5.4), Inches(12.3), Pt(1.5))
ln2.fill.solid(); ln2.fill.fore_color.rgb = BLUE; ln2.line.fill.background()
tx(s,0.5,5.8,12.3,0.7,"感谢各位老师批评指正!  请各位老师提问",sz=Pt(32),cl=NAVY,b=True,al=PP_ALIGN.CENTER)
pn(s,28); nt(s,"不足和展望真诚具体。答辩老师可能就这些方向提问。")

print("Slides 26-28 done")

out="presentation.pptx"
prs.save(out)
print(f"\nSaved: {out} ({len(prs.slides)} slides) - based on template")
